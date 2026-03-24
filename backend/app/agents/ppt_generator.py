import asyncio
import logging
import os
import tempfile
import urllib.request
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

from app.config import settings
from app.models.schemas import SlideContent

logger = logging.getLogger(__name__)

# ── Modern color palette ─────────────────────────────────────────────
COLOR_PRIMARY = RGBColor(0x0F, 0x17, 0x2A)      # Deep navy
COLOR_PRIMARY_L = RGBColor(0x1E, 0x29, 0x3B)    # Lighter navy
COLOR_ACCENT = RGBColor(0x38, 0x8A, 0xE2)       # Bright blue
COLOR_ACCENT_DARK = RGBColor(0x1B, 0x5E, 0xB8)  # Darker blue
COLOR_HIGHLIGHT = RGBColor(0xF0, 0xA8, 0x30)     # Warm gold
COLOR_TEXT = RGBColor(0x2D, 0x2D, 0x2D)          # Near-black
COLOR_TEXT_LIGHT = RGBColor(0x6B, 0x6B, 0x6B)    # Medium gray
COLOR_LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFC)      # Off-white
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_QUOTE_BG = RGBColor(0xEE, 0xF3, 0xFD)      # Soft blue bg
COLOR_GREEN = RGBColor(0x27, 0xAE, 0x60)          # Green accent
COLOR_DIVIDER = RGBColor(0xDD, 0xDD, 0xDD)        # Divider gray

FONT_TITLE = "Microsoft YaHei"
FONT_BODY = "Microsoft YaHei Light"
FONT_FALLBACK = "Arial"

SLIDE_W = Inches(13.333)  # 16:9 widescreen
SLIDE_H = Inches(7.5)


# ── Helpers ──────────────────────────────────────────────────────────

def _set_font(run, size, bold=False, color=COLOR_TEXT, font_name=FONT_TITLE):
    """Set font properties on a text run with East Asian fallback."""
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", font_name)


def _no_border(shape):
    shape.line.fill.background()


def _set_transparency(shape, pct: int):
    """Set fill transparency (0-100) via XML alpha element."""
    spPr = shape._element.spPr
    solidFill = spPr.find(qn("a:solidFill"))
    if solidFill is not None:
        srgb = solidFill.find(qn("a:srgbClr"))
        if srgb is not None:
            val = str((100 - pct) * 1000)  # alpha: 100000=opaque, 0=transparent
            alpha = srgb.makeelement(qn("a:alpha"), {"val": val})
            srgb.append(alpha)


def _add_page_number(slide, num: int, total: int, color=COLOR_TEXT_LIGHT):
    txBox = slide.shapes.add_textbox(
        SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.5),
        Inches(1), Inches(0.4),
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{num} / {total}"
    _set_font(run, 9, color=color, font_name=FONT_BODY)


def _add_decorative_bar(slide, y, width, height, color):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), y, width, height,
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    _no_border(bar)
    return bar


def _add_top_title_bar(slide, title_text: str):
    """Shared top dark title bar used by content, two_column, timeline slides."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(1.15),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    _no_border(bar)

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.18), Inches(11), Inches(0.8))
    tf = txBox.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    _set_font(run, 24, bold=True, color=COLOR_WHITE)

    _add_decorative_bar(slide, Inches(1.15), Inches(2.5), Pt(4), COLOR_ACCENT)


def _add_bottom_bars(slide):
    """Shared bottom decorative bars: left accent + bottom gold line."""
    _add_decorative_bar(slide, Inches(1.15), Pt(5), SLIDE_H - Inches(1.15), COLOR_ACCENT)
    _add_decorative_bar(slide, SLIDE_H - Pt(4), SLIDE_W, Pt(4), COLOR_HIGHLIGHT)


# ── Image utilities ──────────────────────────────────────────────────

def _download_image(url: str) -> str | None:
    """Download image to temp file. Returns path or None on failure."""
    if not url:
        return None
    try:
        suffix = ".png" if ".png" in url.lower() else ".jpg"
        fd, path = tempfile.mkstemp(suffix=suffix)
        os.close(fd)
        req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
        with urllib.request.urlopen(req, timeout=15) as resp:
            with open(path, "wb") as f:
                f.write(resp.read())
        return path
    except Exception as e:
        logger.warning(f"Failed to download image {url[:80]}: {e}")
        return None


def _add_image_to_slide(slide, img_path: str, left, top, width=None, height=None):
    """Add image to slide, cleaning up temp file afterward."""
    try:
        kwargs = {"left": left, "top": top}
        if width:
            kwargs["width"] = width
        if height:
            kwargs["height"] = height
        slide.shapes.add_picture(img_path, **kwargs)
    except Exception as e:
        logger.warning(f"Failed to add image to slide: {e}")
    finally:
        try:
            os.unlink(img_path)
        except OSError:
            pass


def _add_letter_avatar(slide, name: str, name_cn: str):
    """Add a circular letter-avatar shape."""
    avatar = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.8), Inches(0.7), Inches(1.0), Inches(1.0),
    )
    avatar.fill.solid()
    avatar.fill.fore_color.rgb = COLOR_GREEN
    _no_border(avatar)
    atf = avatar.text_frame
    atf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ap = atf.paragraphs[0]
    ap.alignment = PP_ALIGN.CENTER
    ar = ap.add_run()
    ar.text = name_cn[0] if name_cn else (name[0] if name else "?")
    _set_font(ar, 28, bold=True, color=COLOR_WHITE)


# ── Title slide ──────────────────────────────────────────────────────

def _add_title_slide(prs: Presentation, title: str, subtitle: str = "", thumbnail_url: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLOR_PRIMARY

    # Right-side decorative block
    deco = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, SLIDE_W - Inches(4.5), Inches(0), Inches(4.5), SLIDE_H,
    )
    deco.fill.solid()
    deco.fill.fore_color.rgb = COLOR_PRIMARY_L
    _no_border(deco)
    _set_transparency(deco, 60)

    # Thumbnail image on the right side (if available)
    if thumbnail_url:
        img_path = _download_image(thumbnail_url)
        if img_path:
            _add_image_to_slide(
                slide, img_path,
                left=SLIDE_W - Inches(4.3), top=Inches(0.2),
                width=Inches(4.0), height=Inches(7.1),
            )
            # Dark overlay for text readability
            overlay = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                SLIDE_W - Inches(4.3), Inches(0), Inches(4.3), SLIDE_H,
            )
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = COLOR_PRIMARY
            _no_border(overlay)
            _set_transparency(overlay, 45)

    # Gold accent line
    accent_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(3.0), Inches(1.5), Pt(4),
    )
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = COLOR_HIGHLIGHT
    _no_border(accent_line)

    # Title
    txBox = slide.shapes.add_textbox(Inches(1.2), Inches(3.3), Inches(7.5), Inches(1.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    _set_font(run, 40, bold=True, color=COLOR_WHITE)

    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        p2.space_before = Pt(16)
        run2 = p2.add_run()
        run2.text = subtitle
        _set_font(run2, 20, color=RGBColor(0xA0, 0xAE, 0xC4), font_name=FONT_BODY)


# ── Content slide ────────────────────────────────────────────────────

def _add_content_slide(prs: Presentation, slide_content: SlideContent, page_num: int = 0, total_pages: int = 0):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLOR_WHITE

    _add_top_title_bar(slide, slide_content.title)

    if slide_content.bullet_points:
        for i, point in enumerate(slide_content.bullet_points):
            y_offset = Inches(1.6) + Inches(i * 0.62)

            # Numbered circle
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(0.8), y_offset + Pt(2), Inches(0.32), Inches(0.32),
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = COLOR_ACCENT
            _no_border(circle)
            ctf = circle.text_frame
            ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
            cp = ctf.paragraphs[0]
            cp.alignment = PP_ALIGN.CENTER
            cr = cp.add_run()
            cr.text = str(i + 1)
            _set_font(cr, 11, bold=True, color=COLOR_WHITE, font_name=FONT_FALLBACK)

            # Text
            txBox2 = slide.shapes.add_textbox(Inches(1.3), y_offset, Inches(11), Inches(0.55))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
            p2 = tf2.paragraphs[0]
            run2 = p2.add_run()
            run2.text = point
            _set_font(run2, 16, color=COLOR_TEXT, font_name=FONT_BODY)

    _add_bottom_bars(slide)
    if page_num:
        _add_page_number(slide, page_num, total_pages)


# ── Quote slide ──────────────────────────────────────────────────────

def _add_quote_slide(prs: Presentation, slide_content: SlideContent, page_num: int = 0, total_pages: int = 0):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLOR_QUOTE_BG

    # Left decorative bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(1.8), Pt(8), Inches(3.8),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_ACCENT
    _no_border(bar)

    # Large quote mark
    txBox_mark = slide.shapes.add_textbox(Inches(1.4), Inches(1.2), Inches(1.5), Inches(1.2))
    tf_mark = txBox_mark.text_frame
    p_mark = tf_mark.paragraphs[0]
    run_mark = p_mark.add_run()
    run_mark.text = "\u201C"
    _set_font(run_mark, 80, bold=True, color=COLOR_ACCENT)

    # Quote text
    txBox = slide.shapes.add_textbox(Inches(1.8), Inches(2.5), Inches(9.5), Inches(3))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = Pt(32)
    run = p.add_run()
    run.text = slide_content.quote or slide_content.title
    _set_font(run, 24, color=COLOR_PRIMARY, font_name=FONT_BODY)

    if slide_content.speaker:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        p2.space_before = Pt(28)
        run_dash = p2.add_run()
        run_dash.text = "\u2014\u2014  "
        _set_font(run_dash, 16, color=COLOR_HIGHLIGHT)
        run2 = p2.add_run()
        run2.text = slide_content.speaker
        _set_font(run2, 18, bold=True, color=COLOR_ACCENT_DARK)

    _add_decorative_bar(slide, SLIDE_H - Pt(4), SLIDE_W, Pt(4), COLOR_ACCENT)
    if page_num:
        _add_page_number(slide, page_num, total_pages, color=COLOR_ACCENT_DARK)


# ── Summary slide ────────────────────────────────────────────────────

def _add_summary_slide(prs: Presentation, slide_content: SlideContent, page_num: int = 0, total_pages: int = 0):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLOR_PRIMARY

    # Top-right gold bar
    deco = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, SLIDE_W - Inches(3), Inches(0), Inches(3), Pt(6),
    )
    deco.fill.solid()
    deco.fill.fore_color.rgb = COLOR_HIGHLIGHT
    _no_border(deco)

    # Gold accent under title
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(1.4), Inches(2), Pt(4),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLOR_HIGHLIGHT
    _no_border(accent)

    # Title
    txBox = slide.shapes.add_textbox(Inches(1.2), Inches(0.5), Inches(10), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = slide_content.title
    _set_font(run, 30, bold=True, color=COLOR_WHITE)

    if slide_content.bullet_points:
        for i, point in enumerate(slide_content.bullet_points):
            y_offset = Inches(1.9) + Inches(i * 0.58)
            check_box = slide.shapes.add_textbox(Inches(1.2), y_offset, Inches(0.4), Inches(0.4))
            chtf = check_box.text_frame
            chp = chtf.paragraphs[0]
            chr_ = chp.add_run()
            chr_.text = "\u2713"
            _set_font(chr_, 16, bold=True, color=COLOR_HIGHLIGHT, font_name=FONT_FALLBACK)

            txBox2 = slide.shapes.add_textbox(Inches(1.8), y_offset, Inches(10), Inches(0.5))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            run2 = p2.add_run()
            run2.text = point
            _set_font(run2, 17, color=COLOR_WHITE, font_name=FONT_BODY)

    if page_num:
        _add_page_number(slide, page_num, total_pages, color=RGBColor(0x70, 0x80, 0x9A))


# ── Two-column slide (NEW) ───────────────────────────────────────────

def _add_two_column_slide(prs: Presentation, slide_content: SlideContent, page_num: int = 0, total_pages: int = 0):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLOR_WHITE

    _add_top_title_bar(slide, slide_content.title)

    mid_x = SLIDE_W / 2

    # Left column header
    left_hdr = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5))
    ltf = left_hdr.text_frame
    lp = ltf.paragraphs[0]
    lr = lp.add_run()
    lr.text = slide_content.left_title or "A"
    _set_font(lr, 20, bold=True, color=COLOR_ACCENT)

    # Left column points
    for i, point in enumerate(slide_content.left_points[:6]):
        y = Inches(2.2) + Inches(i * 0.6)
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.8), y + Pt(4), Inches(0.15), Inches(0.15),
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = COLOR_ACCENT
        _no_border(dot)

        txBox = slide.shapes.add_textbox(Inches(1.1), y, Inches(5.2), Inches(0.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = point
        _set_font(run, 14, color=COLOR_TEXT, font_name=FONT_BODY)

    # Center divider
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, mid_x - Pt(1), Inches(1.5), Pt(2), Inches(5.5),
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = COLOR_DIVIDER
    _no_border(divider)

    # Right column header
    right_x = Inches(7.0)
    right_hdr = slide.shapes.add_textbox(right_x, Inches(1.5), Inches(5.5), Inches(0.5))
    rtf = right_hdr.text_frame
    rp = rtf.paragraphs[0]
    rr = rp.add_run()
    rr.text = slide_content.right_title or "B"
    _set_font(rr, 20, bold=True, color=COLOR_ACCENT_DARK)

    # Right column points
    for i, point in enumerate(slide_content.right_points[:6]):
        y = Inches(2.2) + Inches(i * 0.6)
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, right_x, y + Pt(4), Inches(0.15), Inches(0.15),
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = COLOR_ACCENT_DARK
        _no_border(dot)

        txBox = slide.shapes.add_textbox(right_x + Inches(0.3), y, Inches(5.2), Inches(0.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = point
        _set_font(run, 14, color=COLOR_TEXT, font_name=FONT_BODY)

    _add_bottom_bars(slide)
    if page_num:
        _add_page_number(slide, page_num, total_pages)


# ── Highlight slide (NEW) ────────────────────────────────────────────

def _add_highlight_slide(prs: Presentation, slide_content: SlideContent, page_num: int = 0, total_pages: int = 0):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLOR_ACCENT

    # Small title at top
    txBox_title = slide.shapes.add_textbox(Inches(1.5), Inches(0.8), Inches(10), Inches(0.6))
    tf_title = txBox_title.text_frame
    p_title = tf_title.paragraphs[0]
    run_title = p_title.add_run()
    run_title.text = slide_content.title or ""
    _set_font(run_title, 20, color=COLOR_WHITE, font_name=FONT_BODY)

    # Gold decorative line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(1.6), Inches(2), Pt(4),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_HIGHLIGHT
    _no_border(line)

    # Large centered highlight text
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.3), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = slide_content.highlight_text or slide_content.title
    _set_font(run, 44, bold=True, color=COLOR_WHITE)

    # Supplementary points below
    if slide_content.bullet_points:
        for i, point in enumerate(slide_content.bullet_points[:3]):
            y = Inches(5.2) + Inches(i * 0.5)
            txBox2 = slide.shapes.add_textbox(Inches(2), y, Inches(9.3), Inches(0.4))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.alignment = PP_ALIGN.CENTER
            run2 = p2.add_run()
            run2.text = point
            _set_font(run2, 16, color=RGBColor(0xCC, 0xDD, 0xF5), font_name=FONT_BODY)

    if page_num:
        _add_page_number(slide, page_num, total_pages, color=RGBColor(0xCC, 0xDD, 0xF5))


# ── Timeline slide (NEW) ─────────────────────────────────────────────

def _add_timeline_slide(prs: Presentation, slide_content: SlideContent, page_num: int = 0, total_pages: int = 0):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLOR_LIGHT_BG

    _add_top_title_bar(slide, slide_content.title)

    # Vertical timeline line
    timeline_x = Inches(2.5)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, timeline_x, Inches(1.5), Pt(3), Inches(5.5),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT
    _no_border(line)

    for i, point in enumerate(slide_content.bullet_points[:8]):
        y = Inches(1.7) + Inches(i * 0.68)

        # Circle node
        node_color = COLOR_ACCENT if i % 2 == 0 else COLOR_HIGHLIGHT
        node = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, timeline_x - Inches(0.12), y, Inches(0.28), Inches(0.28),
        )
        node.fill.solid()
        node.fill.fore_color.rgb = node_color
        _no_border(node)

        # Event text
        txBox = slide.shapes.add_textbox(timeline_x + Inches(0.5), y - Pt(2), Inches(9.5), Inches(0.55))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = point
        _set_font(run, 15, color=COLOR_TEXT, font_name=FONT_BODY)

    _add_bottom_bars(slide)
    if page_num:
        _add_page_number(slide, page_num, total_pages)


# ── Person slide ─────────────────────────────────────────────────────

def _add_person_slide(prs: Presentation, name: str, name_cn: str, context: str,
                      page_num: int = 0, total_pages: int = 0, thumbnail_url: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLOR_LIGHT_BG

    # Top accent bars
    bar1 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.12),
    )
    bar1.fill.solid()
    bar1.fill.fore_color.rgb = COLOR_PRIMARY
    _no_border(bar1)

    bar2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.12), Inches(4), Pt(4),
    )
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = COLOR_GREEN
    _no_border(bar2)

    # Avatar: real image or letter fallback
    if thumbnail_url:
        img_path = _download_image(thumbnail_url)
        if img_path:
            _add_image_to_slide(
                slide, img_path,
                left=Inches(0.8), top=Inches(0.7),
                width=Inches(1.0), height=Inches(1.0),
            )
        else:
            _add_letter_avatar(slide, name, name_cn)
    else:
        _add_letter_avatar(slide, name, name_cn)

    # Person name (Chinese)
    txBox = slide.shapes.add_textbox(Inches(2.1), Inches(0.7), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = name_cn or name
    _set_font(run, 28, bold=True, color=COLOR_PRIMARY)

    # Person name (English)
    if name_cn and name:
        txBox_en = slide.shapes.add_textbox(Inches(2.1), Inches(1.25), Inches(9), Inches(0.4))
        tf_en = txBox_en.text_frame
        p_en = tf_en.paragraphs[0]
        run_en = p_en.add_run()
        run_en.text = name
        _set_font(run_en, 16, color=COLOR_TEXT_LIGHT, font_name=FONT_BODY)

    # Divider
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.0), Inches(11.7), Pt(1.5),
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = COLOR_DIVIDER
    _no_border(divider)

    # Context card
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.3), Inches(11.7), Inches(4.5),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_WHITE
    card.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    card.line.width = Pt(1)

    txBox2 = slide.shapes.add_textbox(Inches(1.2), Inches(2.6), Inches(10.9), Inches(4.0))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.line_spacing = Pt(26)
    run2 = p2.add_run()
    run2.text = context
    _set_font(run2, 16, color=COLOR_TEXT, font_name=FONT_BODY)

    _add_decorative_bar(slide, SLIDE_H - Pt(4), SLIDE_W, Pt(4), COLOR_GREEN)
    if page_num:
        _add_page_number(slide, page_num, total_pages)


# ── TOC slide ────────────────────────────────────────────────────────

def _add_toc_slide(prs: Presentation, items: list[str], page_num: int = 0, total_pages: int = 0):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLOR_WHITE

    _add_top_title_bar(slide, "\u76EE\u5F55")

    col_break = 7
    for i, item in enumerate(items[:14]):
        col = 0 if i < col_break else 1
        row = i if i < col_break else i - col_break
        x_base = Inches(1.0) + col * Inches(6.2)
        y_offset = Inches(1.7) + Inches(row * 0.7)

        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x_base, y_offset + Pt(2), Inches(0.38), Inches(0.30),
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = COLOR_ACCENT if col == 0 else COLOR_ACCENT_DARK
        _no_border(badge)
        btf = badge.text_frame
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = btf.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run()
        br.text = str(i + 1)
        _set_font(br, 11, bold=True, color=COLOR_WHITE, font_name=FONT_FALLBACK)

        txBox2 = slide.shapes.add_textbox(x_base + Inches(0.52), y_offset, Inches(5.3), Inches(0.4))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = item
        _set_font(run2, 15, color=COLOR_TEXT, font_name=FONT_BODY)

    _add_bottom_bars(slide)
    if page_num:
        _add_page_number(slide, page_num, total_pages)


# ── Section divider slide ────────────────────────────────────────────

def _add_section_slide(prs: Presentation, title: str, subtitle: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLOR_ACCENT

    # Gold bar
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.55), Inches(1.8), Pt(4),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_HIGHLIGHT
    _no_border(line)

    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.8), Inches(10), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    _set_font(run, 36, bold=True, color=COLOR_WHITE)

    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        p2.space_before = Pt(12)
        run2 = p2.add_run()
        run2.text = subtitle
        _set_font(run2, 18, color=RGBColor(0xCC, 0xDD, 0xF5), font_name=FONT_BODY)


# ── End slide ────────────────────────────────────────────────────────

def _add_end_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLOR_PRIMARY

    deco = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, SLIDE_W - Inches(5), Inches(0), Inches(5), SLIDE_H,
    )
    deco.fill.solid()
    deco.fill.fore_color.rgb = COLOR_PRIMARY_L
    _no_border(deco)
    _set_transparency(deco, 70)

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(5.2), Inches(3.0), Inches(3), Pt(4),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_HIGHLIGHT
    _no_border(line)

    txBox = slide.shapes.add_textbox(Inches(2), Inches(3.3), Inches(9.3), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "\u8C22\u8C22\u89C2\u770B"
    _set_font(run, 42, bold=True, color=COLOR_WHITE)

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(12)
    run2 = p2.add_run()
    run2.text = "Thank You"
    _set_font(run2, 20, color=RGBColor(0x8A, 0x9A, 0xB4), font_name=FONT_BODY)


# ── Main generation logic ────────────────────────────────────────────

def _generate_ppt_sync(
    video_id: str,
    title: str,
    title_cn: str,
    slides: list[SlideContent],
    mentioned_people: list,
    thumbnail_url: str = "",
) -> str:
    """Synchronous PPT generation (runs in thread pool)."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    content_count = len(slides)
    person_count = len(mentioned_people)
    total_pages = 1 + 1 + content_count + (1 + person_count if person_count else 0) + 1
    page = 0

    # 1. Title slide with thumbnail
    _add_title_slide(prs, title_cn or title, title if title_cn else "", thumbnail_url=thumbnail_url)
    page += 1

    # 2. Table of contents
    toc_items = [
        s.title for s in slides
        if s.slide_type in ("content", "summary", "two_column", "highlight", "timeline") and s.title
    ][:14]
    page += 1
    _add_toc_slide(prs, toc_items, page_num=page, total_pages=total_pages)

    # 3. Content slides — dispatch by type
    for slide_content in slides:
        page += 1
        st = slide_content.slide_type

        if st == "title":
            _add_title_slide(
                prs, slide_content.title,
                slide_content.bullet_points[0] if slide_content.bullet_points else "",
            )
        elif st == "section_title":
            _add_section_slide(
                prs, slide_content.title,
                slide_content.bullet_points[0] if slide_content.bullet_points else "",
            )
        elif st == "quote":
            _add_quote_slide(prs, slide_content, page_num=page, total_pages=total_pages)
        elif st == "summary":
            _add_summary_slide(prs, slide_content, page_num=page, total_pages=total_pages)
        elif st == "two_column":
            _add_two_column_slide(prs, slide_content, page_num=page, total_pages=total_pages)
        elif st == "highlight":
            _add_highlight_slide(prs, slide_content, page_num=page, total_pages=total_pages)
        elif st == "timeline":
            _add_timeline_slide(prs, slide_content, page_num=page, total_pages=total_pages)
        else:
            _add_content_slide(prs, slide_content, page_num=page, total_pages=total_pages)

    # 4. Person slides
    if mentioned_people:
        page += 1
        _add_section_slide(prs, "\u63D0\u53CA\u4EBA\u7269", "\u4EE5\u4E0B\u4E3A\u89C6\u9891\u4E2D\u63D0\u53CA\u7684\u91CD\u8981\u4EBA\u7269")

        for person in mentioned_people:
            page += 1
            name = person.name if hasattr(person, "name") else person.get("name", "")
            name_cn = person.name_cn if hasattr(person, "name_cn") else person.get("name_cn", "")
            context = person.context if hasattr(person, "context") else person.get("context", "")
            thumb = person.thumbnail_url if hasattr(person, "thumbnail_url") else person.get("thumbnail_url", "")
            _add_person_slide(prs, name, name_cn, context,
                              page_num=page, total_pages=total_pages, thumbnail_url=thumb)

    # 5. End slide
    _add_end_slide(prs)

    # Save
    settings.output_dir.mkdir(parents=True, exist_ok=True)
    safe_title = "".join(c for c in video_id if c.isalnum() or c in "-_")
    filename = f"{safe_title}.pptx"
    filepath = settings.output_dir / filename
    prs.save(str(filepath))

    logger.info(f"PPT saved: {filepath}")
    return filename


async def generate_ppt(
    video_id: str,
    title: str,
    title_cn: str,
    slides: list[SlideContent],
    mentioned_people: list,
    thumbnail_url: str = "",
    progress_callback=None,
) -> str:
    """Generate a PPT file and return the filename."""
    if progress_callback:
        await progress_callback("\u6B63\u5728\u751F\u6210 PPT \u6F14\u793A\u6587\u7A3F...")

    filename = await asyncio.to_thread(
        _generate_ppt_sync, video_id, title, title_cn, slides, mentioned_people, thumbnail_url
    )
    return filename
